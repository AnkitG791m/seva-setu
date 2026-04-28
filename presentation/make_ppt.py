from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── COLOUR PALETTE ──────────────────────────────────────────────────────────
BG       = RGBColor(0x0f, 0x17, 0x2a)   # deep navy
CARD     = RGBColor(0x1e, 0x29, 0x3b)   # card bg
GREEN    = RGBColor(0x10, 0xb9, 0x81)   # brand green
INDIGO   = RGBColor(0x63, 0x66, 0xf1)   # accent purple
WHITE    = RGBColor(0xff, 0xff, 0xff)
MUTED    = RGBColor(0x94, 0xa3, 0xb8)
RED      = RGBColor(0xef, 0x44, 0x44)
YELLOW   = RGBColor(0xf5, 0x9e, 0x0b)

BLANK = prs.slide_layouts[6]   # completely blank layout

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG):
    """Fill slide background."""
    bg_ = slide.background
    fill = bg_.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_rgb=CARD, alpha=None, line_rgb=None, line_w=Pt(1)):
    """Add a filled rectangle (card)."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    """Add a text box."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def tag_box(slide, text, l, t):
    """Small pill tag."""
    r = rect(slide, l, t, len(text)*0.12+0.3, 0.33, fill_rgb=RGBColor(0x0d, 0x2a, 0x1e),
             line_rgb=GREEN, line_w=Pt(0.75))
    txt(slide, text, l+0.08, t+0.04, len(text)*0.12+0.2, 0.3,
        size=9, bold=True, color=GREEN, align=PP_ALIGN.LEFT)

def divider(slide, t_inch, color=GREEN):
    r = rect(slide, 0.6, t_inch, 1.2, 0.045, fill_rgb=color)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — HERO
# ════════════════════════════════════════════════════════════════════════════
s1 = add_slide(); bg(s1)

# Gradient accent blobs (simulated with coloured rects, low opacity not possible — use dark tint)
rect(s1, 0, 0, 5, 7.5, fill_rgb=RGBColor(0x0b, 0x20, 0x1c))   # left tinted block
rect(s1, 5, 0, 8.33, 7.5, fill_rgb=BG)

# Green left accent bar
rect(s1, 0, 0, 0.08, 7.5, fill_rgb=GREEN)

txt(s1, "🌿  HACKATHON 2026 · PROTOTYPE SUBMISSION",
    0.6, 0.35, 7, 0.4, size=9, bold=True, color=GREEN)

txt(s1, "SevaSetu",
    0.6, 0.85, 7, 1.4, size=72, bold=True, color=WHITE)

txt(s1, "Bridge to Hope",
    0.65, 2.0, 7, 0.5, size=16, bold=False, color=MUTED, italic=True)

divider(s1, 2.65)

txt(s1, (
    "An AI-powered platform using Google Gemini and real-time\n"
    "geospatial intelligence to automate crisis triage, prioritize\n"
    "community needs, and coordinate volunteers — ensuring the\n"
    "right help reaches the right place at the right time."
), 0.6, 2.8, 6.8, 2.0, size=14, color=MUTED)

# Badges row
badges = ["Google Cloud Run", "Gemini AI 1.5 Flash", "Google Maps Platform", "Firebase Auth"]
bx = 0.6
for b in badges:
    w = len(b)*0.11 + 0.4
    rect(s1, bx, 5.0, w, 0.38, fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=GREEN, line_w=Pt(0.5))
    txt(s1, "●  " + b, bx+0.1, 5.02, w-0.1, 0.34, size=9, color=WHITE, bold=False)
    bx += w + 0.15

# Right side stat cards
stats = [("⚡ AI", "Urgency Scoring"), ("🌍 Maps", "Geospatial Heatmap"), ("☁️ GCP", "Cloud Native")]
sy = 0.8
for ico, lbl in stats:
    r = rect(s1, 10.2, sy, 2.7, 1.4, fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=GREEN, line_w=Pt(0.5))
    txt(s1, ico, 10.2, sy+0.15, 2.7, 0.7, size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s1, lbl, 10.2, sy+0.85, 2.7, 0.4, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    sy += 1.6

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s2 = add_slide(); bg(s2)
rect(s2, 0, 0, 0.08, 7.5, fill_rgb=RED)

txt(s2, "🔴  THE PROBLEM", 0.6, 0.3, 5, 0.4, size=9, bold=True, color=RED)
txt(s2, "Crisis Response is Broken", 0.6, 0.7, 10, 1.0, size=38, bold=True, color=WHITE)
divider(s2, 1.75, RED)

problems = [
    ("⏳", "Delayed Triage",
     "Authorities manually analyse reports, wasting the critical\n'Golden Hour' while situations worsen on the ground."),
    ("📍", "No Location Intelligence",
     "Volunteers are deployed arbitrarily — without knowing\nwhich zones are most critically affected in real-time."),
    ("🗂️", "Resource Wastage",
     "Duplicate efforts in some areas while high-urgency\nneeds in other zones go completely ignored."),
]
px = 0.6
for ico, title, body in problems:
    r = rect(s2, px, 2.1, 4.0, 3.4, fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=RED, line_w=Pt(0.5))
    rect(s2, px, 2.1, 4.0, 0.06, fill_rgb=RED)   # top accent bar
    txt(s2, ico, px+0.2, 2.25, 0.7, 0.7, size=28)
    txt(s2, title, px+0.2, 2.9, 3.6, 0.5, size=14, bold=True, color=RGBColor(0xfc, 0xa5, 0xa5))
    txt(s2, body, px+0.2, 3.4, 3.5, 1.6, size=11, color=MUTED)
    px += 4.2

# Bottom quote box
rect(s2, 0.6, 5.7, 12.1, 1.25, fill_rgb=RGBColor(0x1a, 0x10, 0x10), line_rgb=RED, line_w=Pt(0.5))
txt(s2, (
    '"During the 2023 Flood Response in Assam, 40% of NGO resources were misdirected '
    'due to lack of a central prioritization system — costing 72+ hours of critical response time."'
), 0.9, 5.85, 11.5, 1.0, size=11, italic=True, color=RGBColor(0xfc, 0xa5, 0xa5))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ════════════════════════════════════════════════════════════════════════════
s3 = add_slide(); bg(s3)
rect(s3, 0, 0, 0.08, 7.5, fill_rgb=GREEN)

txt(s3, "💡  OUR SOLUTION", 0.6, 0.3, 5, 0.4, size=9, bold=True, color=GREEN)
txt(s3, "SevaSetu: End-to-End AI Coordination", 0.6, 0.7, 11, 1.0, size=32, bold=True, color=WHITE)
divider(s3, 1.75)

features = [
    ("🤖", "AI-Driven Urgency Scoring",
     "Google Gemini 1.5 Flash analyses every report\nand assigns an urgency score 0–100 instantly.\nMedical, water, food, shelter — prioritised in seconds."),
    ("🗺️", "Geospatial Crisis Heatmap",
     "Google Maps displays a real-time heatmap of active\nneed reports so coordinators see exactly where to\ndeploy volunteers for maximum impact."),
    ("📸", "Citizen-Powered Reporting",
     "Community members submit geo-tagged reports with\nphotos (stored in GCS) via a mobile-friendly UI with\nPlaces Autocomplete for accurate location input."),
    ("🔔", "Instant Volunteer Alerts",
     "Firebase Cloud Messaging (FCM) pushes real-time task\nassignments to volunteers the moment a high-urgency\nneed is identified in their area."),
]
positions = [(0.6, 2.0), (6.77, 2.0), (0.6, 4.7), (6.77, 4.7)]
for (l, t), (ico, title, body) in zip(positions, features):
    rect(s3, l, t, 5.9, 2.4, fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=GREEN, line_w=Pt(0.5))
    txt(s3, ico, l+0.2, t+0.15, 0.8, 0.7, size=24)
    txt(s3, title, l+0.2, t+0.8, 5.4, 0.5, size=13, bold=True, color=WHITE)
    txt(s3, body, l+0.2, t+1.28, 5.4, 1.0, size=10, color=MUTED)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
s4 = add_slide(); bg(s4)
rect(s4, 0, 0, 0.08, 7.5, fill_rgb=INDIGO)

txt(s4, "🏗️  TECHNICAL ARCHITECTURE", 0.6, 0.3, 6, 0.4, size=9, bold=True, color=INDIGO)
txt(s4, "Built on Google Cloud", 0.6, 0.7, 9, 0.9, size=34, bold=True, color=WHITE)
divider(s4, 1.65, INDIGO)

# Flow row 1
nodes_r1 = [
    ("📱", "Community", "React + Vite"),
    ("🔐", "Firebase Auth", "Google Sign-In"),
    ("☁️", "Cloud Run API", "Node.js / Express"),
    ("🤖", "Gemini AI", "1.5 Flash Model"),
]
nx = 0.5
for ico, name, sub in nodes_r1:
    r = rect(s4, nx, 2.0, 2.7, 1.5, fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=GREEN, line_w=Pt(0.75))
    txt(s4, ico,  nx, 2.1,  2.7, 0.6, size=24, align=PP_ALIGN.CENTER)
    txt(s4, name, nx, 2.65, 2.7, 0.4, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s4, sub,  nx, 3.02, 2.7, 0.4, size=9,  color=MUTED, align=PP_ALIGN.CENTER)
    nx += 2.7
    if nx < 11.3:
        txt(s4, "→", nx-0.15, 2.55, 0.5, 0.5, size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        nx += 0.2

# Flow row 2
nodes_r2 = [
    ("🗄️", "Prisma ORM",     "Structured DB"),
    ("🪣", "Cloud Storage",   "Photo Uploads"),
    ("📍", "Maps Platform",   "Heatmap + Places"),
    ("🔔", "FCM Alerts",      "Push Notifications"),
]
nx2 = 0.5
for ico, name, sub in nodes_r2:
    rect(s4, nx2, 3.9, 2.9, 1.4, fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=INDIGO, line_w=Pt(0.5))
    txt(s4, ico,  nx2, 3.98, 2.9, 0.55, size=22, align=PP_ALIGN.CENTER)
    txt(s4, name, nx2, 4.5,  2.9, 0.38, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s4, sub,  nx2, 4.85, 2.9, 0.35, size=9,  color=MUTED, align=PP_ALIGN.CENTER)
    nx2 += 3.1

# GCP badges
gcp_tags = ["🌏 Region: asia-south1 (Mumbai)", "🐳 Containerized via Docker",
            "🏗️ Cloud Build CI/CD", "🔥 Firebase Hosting (Frontend)"]
bx4 = 0.5
for b in gcp_tags:
    w = len(b)*0.1 + 0.5
    rect(s4, bx4, 5.75, w, 0.42, fill_rgb=RGBColor(0x13, 0x13, 0x35), line_rgb=INDIGO, line_w=Pt(0.5))
    txt(s4, b, bx4+0.1, 5.8, w, 0.35, size=9, color=WHITE)
    bx4 += w + 0.2

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — GOOGLE SERVICES
# ════════════════════════════════════════════════════════════════════════════
s5 = add_slide(); bg(s5)
rect(s5, 0, 0, 0.08, 7.5, fill_rgb=GREEN)

txt(s5, "🔧  GOOGLE SERVICES INTEGRATION", 0.6, 0.3, 7, 0.4, size=9, bold=True, color=GREEN)
txt(s5, "Powered by Google Ecosystem", 0.6, 0.7, 10, 0.9, size=34, bold=True, color=WHITE)
divider(s5, 1.65)

services = [
    ("🤖", "Google Gemini AI 1.5 Flash",
     "Auto-analyses each report in Hindi/English. Assigns urgency\nscore, priority label (CRITICAL / HIGH / MEDIUM / LOW)\nand a one-line reason for coordinators.",
     "PRIMARY AI MODEL"),
    ("🗺️", "Google Maps Platform",
     "Interactive crisis heatmap, marker clusters, and Places\nAutocomplete for accurate geo-tagged location input\nduring report submission by community members.",
     "MAPS JS API + PLACES API"),
    ("🔐", "Firebase Authentication",
     "Secure Google Sign-In and Phone OTP authentication\nfor all users — community members, volunteers,\nand NGO coordinators. Backed by Firebase Admin SDK.",
     "FIREBASE AUTH + FCM"),
    ("☁️", "Google Cloud Run + Cloud Build",
     "Fully containerized, auto-scaling backend in Mumbai\n(asia-south1) delivering low-latency response across India.\nBuilt with CI/CD via Cloud Build pipeline.",
     "CLOUD RUN + CLOUD BUILD"),
]
sx = 0.5
for ico, title, body, tag in services:
    r = rect(s5, sx, 2.0, 5.8, 3.3, fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=GREEN, line_w=Pt(0.5))
    txt(s5, ico,   sx+0.2, 2.1,  0.8, 0.8, size=28)
    txt(s5, title, sx+0.2, 2.85, 5.2, 0.5, size=13, bold=True, color=GREEN)
    txt(s5, body,  sx+0.2, 3.35, 5.2, 1.4, size=10, color=MUTED)
    # tag pill
    w_tag = len(tag)*0.09 + 0.25
    rect(s5, sx+0.2, 4.85, w_tag, 0.3, fill_rgb=RGBColor(0x0d, 0x2a, 0x1e), line_rgb=GREEN, line_w=Pt(0.3))
    txt(s5, tag, sx+0.28, 4.87, w_tag, 0.28, size=8, bold=True, color=GREEN)
    sx += 6.1

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LIVE DEMO
# ════════════════════════════════════════════════════════════════════════════
s6 = add_slide(); bg(s6)
rect(s6, 0, 0, 0.08, 7.5, fill_rgb=INDIGO)

txt(s6, "🚀  LIVE PROTOTYPE", 0.6, 0.3, 5, 0.4, size=9, bold=True, color=INDIGO)
txt(s6, "Fully Deployed & Production Ready", 0.6, 0.7, 11, 0.9, size=32, bold=True, color=WHITE)
divider(s6, 1.65, INDIGO)

links = [
    ("🌐", "LIVE FRONTEND", "https://seva-setu-74bcc.web.app"),
    ("⚡", "LIVE API (Cloud Run)", "https://sevasetu-api-1032768844799.asia-south1.run.app"),
    ("💻", "GITHUB REPOSITORY", "https://github.com/AnkitG791m/seva-setu"),
]
ly = 2.0
for ico, lbl, url in links:
    rect(s6, 0.6, ly, 7.3, 0.75, fill_rgb=RGBColor(0x0f, 0x1c, 0x35), line_rgb=INDIGO, line_w=Pt(0.5))
    txt(s6, ico, 0.75, ly+0.12, 0.4, 0.55, size=16)
    txt(s6, lbl, 1.25, ly+0.05, 2.5, 0.35, size=8, bold=True, color=INDIGO)
    txt(s6, url, 1.25, ly+0.35, 6.5, 0.35, size=10, color=RGBColor(0x93, 0xc5, 0xfd))
    ly += 0.9

# Flow steps
steps = [
    ("1", "Google Login → Role Selection"),
    ("2", "Community member submits Need report with photo + location"),
    ("3", "Gemini AI scores it instantly (e.g. CRITICAL: 87/100)"),
    ("4", "Maps heatmap updates — coordinator assigns nearest volunteer"),
    ("5", "FCM Push notification sent to volunteer in real-time"),
]
txt(s6, "User Flow Walkthrough", 0.6, 5.05, 6, 0.4, size=13, bold=True, color=WHITE)
fy = 5.45
for num, step in steps:
    rect(s6, 0.6, fy, 0.38, 0.38, fill_rgb=GREEN)
    txt(s6, num, 0.6, fy+0.02, 0.38, 0.36, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s6, step, 1.1, fy+0.05, 6.7, 0.35, size=10, color=MUTED)
    fy += 0.42

# Right side mock card
rect(s6, 8.5, 2.0, 4.2, 5.0, fill_rgb=RGBColor(0x1e, 0x29, 0x3b), line_rgb=INDIGO, line_w=Pt(0.5))
# window chrome dots
rect(s6, 8.65, 2.15, 0.18, 0.18, fill_rgb=RED)
rect(s6, 8.9,  2.15, 0.18, 0.18, fill_rgb=YELLOW)
rect(s6, 9.15, 2.15, 0.18, 0.18, fill_rgb=GREEN)
txt(s6, "SevaSetu App", 8.5, 2.55, 4.2, 0.5, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s6, "🌉",           8.5, 3.05, 4.2, 0.8, size=36, align=PP_ALIGN.CENTER)
tags_mock = [
    (GREEN, "✓ Google Login Active"),
    (RGBColor(0x1e,0x29,0x3b), "🤖 AI: Scoring Reports..."),
    (RED,   "🔴 CRITICAL · Urgency: 91/100"),
]
my = 3.9
for bc, mt in tags_mock:
    rect(s6, 8.75, my, 3.6, 0.45, fill_rgb=RGBColor(0x0f,0x17,0x2a), line_rgb=bc, line_w=Pt(0.5))
    txt(s6, mt, 8.85, my+0.07, 3.4, 0.35, size=10, color=WHITE)
    my += 0.58

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — FUTURE ROADMAP
# ════════════════════════════════════════════════════════════════════════════
s7 = add_slide(); bg(s7)
rect(s7, 0, 0, 0.08, 7.5, fill_rgb=YELLOW)

txt(s7, "🔭  FUTURE ROADMAP", 0.6, 0.3, 5, 0.4, size=9, bold=True, color=YELLOW)
txt(s7, "What's Next for SevaSetu", 0.6, 0.7, 10, 0.9, size=34, bold=True, color=WHITE)
divider(s7, 1.65, YELLOW)

timeline = [
    ("🧠", "Phase 2: Predictive Hotspot Detection",
     "Use Vertex AI to analyse historical crisis data and predict\nvulnerable zones before incidents occur — enabling proactive\nresource pre-positioning before disasters strike."),
    ("📡", "Phase 3: Offline-First Mode",
     "SMS/USSD gateway integration so field workers in remote\nareas with no internet can still submit reports and receive\ntask assignments via basic mobile phones."),
    ("🛸", "Phase 4: Drone Zone Mapping",
     "Automated drone deployment to high-urgency zones identified\nby SevaSetu AI, enabling aerial surveys, supply drops, and\nreal-time visual confirmation of crisis situations."),
]
ty = 2.05
for ico, title, body in timeline:
    rect(s7, 0.6, ty, 0.55, 0.55, fill_rgb=RGBColor(0x1e,0x29,0x3b), line_rgb=GREEN, line_w=Pt(1))
    txt(s7, ico, 0.6, ty+0.03, 0.55, 0.5, size=18, align=PP_ALIGN.CENTER)
    txt(s7, title, 1.35, ty,      5.8, 0.42, size=13, bold=True, color=WHITE)
    txt(s7, body,  1.35, ty+0.42, 5.8, 1.0,  size=10, color=MUTED)
    rect(s7, 0.86, ty+0.55, 0.03, 1.05, fill_rgb=GREEN)   # timeline line
    ty += 1.65

# Scale targets table
txt(s7, "📊 Scale Targets", 8.2, 2.0, 4.5, 0.5, size=14, bold=True, color=GREEN)
scale = [
    ("Reports / day",     "100,000+"),
    ("Active Volunteers", "50,000+"),
    ("Response Time",     "< 15 mins"),
    ("States Covered",    "All 28 States"),
]
sy7 = 2.6
for label, val in scale:
    rect(s7, 8.2, sy7, 4.5, 0.55, fill_rgb=RGBColor(0x1e,0x29,0x3b))
    txt(s7, label, 8.35, sy7+0.1, 3.0, 0.38, size=10, color=MUTED)
    txt(s7, val,   11.0, sy7+0.1, 1.6, 0.38, size=12, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
    sy7 += 0.62

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CLOSING
# ════════════════════════════════════════════════════════════════════════════
s8 = add_slide(); bg(s8)

# Circular glow sim
for r_size, alpha_fill in [(4.5, RGBColor(0x0d,0x22,0x18)), (3.2, RGBColor(0x0e,0x25,0x1b)),
                            (1.8, RGBColor(0x10,0x28,0x1e))]:
    cx, cy = 6.67, 3.75
    rect(s8, cx-r_size/2, cy-r_size/2*0.75, r_size, r_size*0.75, fill_rgb=alpha_fill)

txt(s8, "🌉", 5.2, 1.1, 2.9, 1.2, size=54, align=PP_ALIGN.CENTER)

txt(s8, "Building a Safer India",
    1.5, 2.25, 9.8, 1.3, size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s8, (
    "SevaSetu is more than a platform — it's a commitment to ensuring that\n"
    "no community need goes unheard, no volunteer is misdirected, and every\n"
    "crisis is met with the precision of artificial intelligence."
), 2.0, 3.55, 9.3, 1.5, size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# Link pills
links8 = [("🌐 Live App", "seva-setu-74bcc.web.app"),
           ("💻 GitHub",  "AnkitG791m/seva-setu"),
           ("⚡ API",     "Cloud Run · asia-south1")]
lx8 = 2.4
for lbl, sub in links8:
    rect(s8, lx8, 5.2, 2.8, 0.75, fill_rgb=RGBColor(0x1e,0x29,0x3b), line_rgb=GREEN, line_w=Pt(0.5))
    txt(s8, lbl, lx8, 5.27, 2.8, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s8, sub, lx8, 5.6,  2.8, 0.32, size=8,  color=GREEN, align=PP_ALIGN.CENTER)
    lx8 += 3.0

txt(s8, "● Deployed on Google Cloud     ● Powered by Gemini AI     ● Firebase + Maps Platform",
    1.5, 6.5, 10.3, 0.5, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ── SAVE ────────────────────────────────────────────────────────────────────
out = r"c:\Users\User\Desktop\SevaSetu\presentation\SevaSetu_Prototype_Deck.pptx"
prs.save(out)
print("Saved: " + out)

